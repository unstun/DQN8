"""就地给 PPTX 中所有图片添加虚线边框，不修改其他内容。"""
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE

BORDER_COLOR = RGBColor(0x44, 0x72, 0xC4)  # Steel Blue
BORDER_WIDTH = Pt(1.0)
BORDER_DASH = MSO_LINE_DASH_STYLE.DASH

pptx_path = sys.argv[1] if len(sys.argv) > 1 else (
    "/home/sun/phdproject/dqn/DQN8/paperdqn8.3/media/fig_platform.pptx"
)

prs = Presentation(pptx_path)
count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            shape.line.color.rgb = BORDER_COLOR
            shape.line.width = BORDER_WIDTH
            shape.line.dash_style = BORDER_DASH
            count += 1

prs.save(pptx_path)
print(f"Done: added dashed borders to {count} pictures in {pptx_path}")
