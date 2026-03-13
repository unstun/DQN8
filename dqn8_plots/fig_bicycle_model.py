#!/usr/bin/env python3
"""Generate a publication-quality Bicycle Kinematic Model diagram using python-pptx.

Output: paperdqn8.3/media/fig_bicycle_model.pptx
Style : Matches project PPT style (Arial, clean black lines, blue accents).
"""

import math
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Paths ──────────────────────────────────────────────────────────────
PROJ = Path(__file__).resolve().parent.parent
OUT_PPTX = PROJ / "paperdqn8.3" / "media" / "fig_bicycle_model.pptx"

# ── Colour palette ─────────────────────────────────────────────────────
BLACK = RGBColor(0x00, 0x00, 0x00)
BLUE = RGBColor(0x00, 0x72, 0xC6)       # accent blue
GRAY = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xCC, 0x00, 0x00)

# ── Geometry (all in inches, slide coords: Y↓) ────────────────────────
THETA_DEG = 32          # heading angle
DELTA_DEG = 28          # front-wheel steering angle
THETA = math.radians(THETA_DEG)
DELTA = math.radians(DELTA_DEG)
L = 2.0                 # wheelbase length on slide
WHEEL_LEN = 0.38        # wheel rectangle length
WHEEL_WID = 0.11        # wheel rectangle width
TRACK_W = 0.7           # track width (distance between left/right wheels)

# Rear axle centre (reference point)
RX, RY = 1.55, 3.25

# Helper: inches → EMU
def emu(inches):
    return int(inches * 914400)

def polar(cx, cy, angle, dist):
    """From (cx,cy), move along `angle` (math convention) by `dist`.
    Slide Y is inverted, so subtract sin component."""
    return cx + dist * math.cos(angle), cy - dist * math.sin(angle)

def add_line(slide, x1, y1, x2, y2, color=BLACK, width_pt=1.2, dash=None):
    """Add a straight line shape."""
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        emu(x1), emu(y1), emu(x2), emu(y2),
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    if dash:
        connector.line.dash_style = dash
    return connector

def add_arrow_line(slide, x1, y1, x2, y2, color=BLACK, width_pt=1.2):
    """Line with arrowhead at end."""
    conn = add_line(slide, x1, y1, x2, y2, color, width_pt)
    # Add arrowhead via XML
    ln = conn.line._ln
    tail = ln.makeelement(qn('a:tailEnd'), {
        'type': 'triangle', 'w': 'med', 'len': 'med'
    })
    ln.append(tail)
    return conn

def add_rect(slide, cx, cy, w, h, angle_deg=0, fill=BLACK, line_color=BLACK, line_w=0.8,
             rounded=False):
    """Add a rectangle centred at (cx, cy) with rotation."""
    left = cx - w / 2
    top = cy - h / 2
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, emu(left), emu(top), emu(w), emu(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line_color
    shape.line.width = Pt(line_w)
    shape.rotation = -angle_deg  # PPT rotation is clockwise
    # Reduce corner radius for rounded rectangles
    if rounded:
        shape.adjustments[0] = 0.15  # smaller corner radius
    return shape

def add_text(slide, x, y, text, size_pt=9, bold=False, color=BLACK, italic=False,
             anchor='center'):
    """Add a small text box centred at (x, y)."""
    w, h = 1.0, 0.35
    txBox = slide.shapes.add_textbox(emu(x - w / 2), emu(y - h / 2), emu(w), emu(h))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.name = 'Arial'
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txBox

def add_arc_polyline(slide, cx, cy, radius, start_deg, end_deg, n_pts=30,
                     color=BLACK, width_pt=1.0):
    """Draw an arc as a freeform polyline."""
    angles = [math.radians(start_deg + i * (end_deg - start_deg) / n_pts)
              for i in range(n_pts + 1)]
    pts = [(cx + radius * math.cos(a), cy - radius * math.sin(a)) for a in angles]
    builder = slide.shapes.build_freeform(emu(pts[0][0]), emu(pts[0][1]))
    segments = [(emu(px), emu(py)) for px, py in pts[1:]]
    builder.add_line_segments(segments)
    shape = builder.convert_to_shape()
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)
    shape.fill.background()  # no fill
    return shape

def add_dot(slide, cx, cy, r=0.04, fill=BLACK):
    """Small filled circle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        emu(cx - r), emu(cy - r), emu(2 * r), emu(2 * r),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_dashed_circle(slide, cx, cy, r=0.15, color=GRAY, width_pt=0.6):
    """Dashed circle (e.g. for turning radius hint)."""
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        emu(cx - r), emu(cy - r), emu(2 * r), emu(2 * r),
    )
    shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)
    shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return shape

# ── Build slide ────────────────────────────────────────────────────────
def build():
    from pptx.enum.dml import MSO_LINE_DASH_STYLE

    prs = Presentation()
    prs.slide_width = Inches(5.2)
    prs.slide_height = Inches(4.8)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background white
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    # ── Coordinate axes ────────────────────────────────────────────
    ox, oy = 0.65, 4.0   # origin
    ax_len = 4.2
    ay_len = 3.5

    # X-axis
    add_arrow_line(slide, ox, oy, ox + ax_len, oy, BLACK, 1.4)
    add_text(slide, ox + ax_len + 0.15, oy + 0.02, 'X', size_pt=12, bold=True, italic=True)

    # Y-axis
    add_arrow_line(slide, ox, oy, ox, oy - ay_len, BLACK, 1.4)
    add_text(slide, ox - 0.02, oy - ay_len - 0.2, 'Y', size_pt=12, bold=True, italic=True)

    # Origin label
    add_text(slide, ox - 0.18, oy + 0.18, 'O', size_pt=9, italic=True)

    # ── Key geometry ──────────────────────────────────────────────
    perp = THETA + math.pi / 2  # perpendicular to heading
    half_track = TRACK_W / 2
    fx, fy = polar(RX, RY, THETA, L)  # front axle centre

    # ── Dashed projection lines (draw first, behind everything) ───
    conn3 = add_line(slide, RX, RY, RX, oy, GRAY, 0.5)
    conn3.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    conn4 = add_line(slide, RX, RY, ox, RY, GRAY, 0.5)
    conn4.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    # ── Heading reference line (dashed, from rear axle along +X) ──
    ref_end_x = RX + 1.3
    conn_ref = add_line(slide, RX, RY, ref_end_x, RY, GRAY, 0.8)
    conn_ref.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    # ── Front heading reference (dashed extension of body) ────────
    ext_len = 0.9
    fref_x, fref_y = polar(fx, fy, THETA, ext_len)
    conn_fref = add_line(slide, fx, fy, fref_x, fref_y, GRAY, 0.8)
    conn_fref.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    # ── Vehicle body line (rear → front axle) ─────────────────────
    add_line(slide, RX, RY, fx, fy, BLACK, 2.2)

    # ── Axle lines (perpendicular to heading) ─────────────────────
    # Rear axle
    rl_x, rl_y = polar(RX, RY, perp, half_track)
    rr_x, rr_y = polar(RX, RY, perp, -half_track)
    add_line(slide, rl_x, rl_y, rr_x, rr_y, BLACK, 1.0)

    # Front axle
    fl_x, fl_y = polar(fx, fy, perp, half_track)
    fr_x, fr_y = polar(fx, fy, perp, -half_track)
    add_line(slide, fl_x, fl_y, fr_x, fr_y, BLACK, 1.0)

    # ── Wheels (4 filled rounded rectangles) ──────────────────────
    DARK = RGBColor(0x2A, 0x2A, 0x2A)
    rear_wheel_angle = math.degrees(THETA)
    front_wheel_angle = math.degrees(THETA + DELTA)

    for wx, wy in [(rl_x, rl_y), (rr_x, rr_y)]:
        add_rect(slide, wx, wy, WHEEL_LEN, WHEEL_WID, rear_wheel_angle,
                 fill=DARK, line_color=BLACK, line_w=0.6, rounded=True)
    for wx, wy in [(fl_x, fl_y), (fr_x, fr_y)]:
        add_rect(slide, wx, wy, WHEEL_LEN, WHEEL_WID, front_wheel_angle,
                 fill=DARK, line_color=BLACK, line_w=0.6, rounded=True)

    # ── Rear axle centre dot ──────────────────────────────────────
    add_dot(slide, RX, RY, r=0.045, fill=BLUE)

    # ── Wheelbase dimension "L" with dimension lines ──────────────
    # Standard engineering dimension: two ticks + connecting line + label
    # Place on the LEFT side of body (positive perp direction)
    dim_off = 0.35  # offset from body centreline
    tick_half = 0.08  # half-length of each tick mark

    # Tick at rear axle
    tr_x, tr_y = polar(RX, RY, perp, dim_off - tick_half)
    tr2_x, tr2_y = polar(RX, RY, perp, dim_off + tick_half)
    add_line(slide, tr_x, tr_y, tr2_x, tr2_y, BLACK, 1.0)
    # Tick at front axle
    tf_x, tf_y = polar(fx, fy, perp, dim_off - tick_half)
    tf2_x, tf2_y = polar(fx, fy, perp, dim_off + tick_half)
    add_line(slide, tf_x, tf_y, tf2_x, tf2_y, BLACK, 1.0)
    # Connecting dimension line between ticks (with arrowheads)
    tc_r_x, tc_r_y = polar(RX, RY, perp, dim_off)
    tc_f_x, tc_f_y = polar(fx, fy, perp, dim_off)
    dim_conn = add_line(slide, tc_r_x, tc_r_y, tc_f_x, tc_f_y, BLACK, 0.9)
    # Add arrowheads on both ends
    ln = dim_conn.line._ln
    ln.append(ln.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'sm', 'len': 'sm'}))
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'sm', 'len': 'sm'}))
    # Extension lines from body to dimension line
    ext_r_x, ext_r_y = polar(RX, RY, perp, 0.08)
    add_line(slide, ext_r_x, ext_r_y, tr2_x, tr2_y, GRAY, 0.5)
    ext_f_x, ext_f_y = polar(fx, fy, perp, 0.08)
    add_line(slide, ext_f_x, ext_f_y, tf2_x, tf2_y, GRAY, 0.5)
    # L label at midpoint, offset further out
    mid_x, mid_y = (tc_r_x + tc_f_x) / 2, (tc_r_y + tc_f_y) / 2
    lx, ly = polar(mid_x, mid_y, perp, 0.20)
    add_text(slide, lx, ly, 'L', size_pt=13, bold=True, italic=True, color=BLACK)

    # ── Heading angle θ arc (small, tight to rear axle) ────────────
    arc_r = 0.38
    add_arc_polyline(slide, RX, RY, arc_r, 0, THETA_DEG, n_pts=30, color=BLUE, width_pt=1.3)
    # θ label — place outside arc, biased toward horizontal
    theta_lbl_angle = THETA / 2
    tlx, tly = polar(RX, RY, theta_lbl_angle, arc_r + 0.20)
    add_text(slide, tlx, tly, 'θ', size_pt=12, bold=True, italic=True, color=BLUE)

    # ── Steering angle δ arc ──────────────────────────────────────
    arc_r_d = 0.42
    add_arc_polyline(slide, fx, fy, arc_r_d, THETA_DEG, THETA_DEG + DELTA_DEG,
                     n_pts=25, color=RED, width_pt=1.3)
    # δ label
    delta_mid = THETA + DELTA / 2
    dlx, dly = polar(fx, fy, delta_mid, arc_r_d + 0.22)
    add_text(slide, dlx, dly, 'δ', size_pt=12, bold=True, italic=True, color=RED)

    # ── Velocity arrow v (on RIGHT side, well clear of θ) ────────
    # Offset along body then perpendicular, so it's visually separate
    v_body_off = 0.30   # start a bit ahead along the body from rear axle
    v_perp_off = -0.25  # offset to the right of body
    v_base_x, v_base_y = polar(RX, RY, THETA, v_body_off)
    v_start_x, v_start_y = polar(v_base_x, v_base_y, perp, v_perp_off)
    v_len = 0.80
    v_end_x, v_end_y = polar(v_start_x, v_start_y, THETA, v_len)
    add_arrow_line(slide, v_start_x, v_start_y, v_end_x, v_end_y, BLUE, 2.2)
    # v label — at the tip, offset to right
    v_lbl_x, v_lbl_y = polar(v_end_x, v_end_y, THETA, 0.05)
    v_lbl_x2, v_lbl_y2 = polar(v_lbl_x, v_lbl_y, perp, -0.20)
    add_text(slide, v_lbl_x2, v_lbl_y2, 'v', size_pt=11, bold=True, italic=True, color=BLUE)

    # ── Front steered direction arrow ─────────────────────────────
    steer_len = 0.7
    sx, sy = polar(fx, fy, THETA + DELTA, steer_len)
    add_arrow_line(slide, fx, fy, sx, sy, RED, 1.6)

    # ── Front axle centre dot (drawn last so it's on top) ────────
    add_dot(slide, fx, fy, r=0.035, fill=RED)

    # ── Rear axle label (x, y) ────────────────────────────────────
    # Place below-right of the rear axle dot
    add_text(slide, RX - 0.30, RY + 0.28, '(x, y)', size_pt=9, italic=True, color=BLACK)

    # ── Save ──────────────────────────────────────────────────────
    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"✓ Saved → {OUT_PPTX}")

if __name__ == "__main__":
    build()
