import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# 基本路径设置
base_dir = "/home/sun/phdproject/dqn/DQN8/作图资料i/图片png/设备图"
main_img_path = os.path.join(base_dir, "大车整体图.png")
sub1_path = os.path.join(base_dir, "雷达mid360.png")
sub2_path = os.path.join(base_dir, "工控机.png")
sub3_path = os.path.join(base_dir, "无人小车底盘.png")

# 1. 裁剪大图并保存临时文件
img_main = Image.open(main_img_path)
CROP_BOX = (0, 70, 405, 545)
img_cropped = img_main.crop(CROP_BOX)

# 创建临时图片，供 python-pptx 读取
tmp_main = "/tmp/main_cropped.png"
img_cropped.save(tmp_main)

# 2. 初始化 PPT
prs = Presentation()
prs.slide_width = Inches(5.5)
prs.slide_height = Inches(4.5)

# 使用空白布局
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# 3. 插入裁切后的大图 (位于左侧)
pic_main = slide.shapes.add_picture(tmp_main, Inches(0.2), Inches(0.5), height=Inches(3.5))

# 4. 辅助函数插入右侧小图及说明文字
def add_subimg(img_path, top_in, label_text):
    # 插入子图，固定宽度 0.9 英寸
    pic = slide.shapes.add_picture(img_path, Inches(3.8), Inches(top_in), width=Inches(0.9))

    # 获取图片实际高度以计算文本框位置
    height_in = pic.height / 914400.0  # EMU to inches

    # 添加文本框
    txBox = slide.shapes.add_textbox(Inches(3.2), Inches(top_in + height_in + 0.05), Inches(2.1), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    # 取消内置边距，防止字太丑折行
    txBox.margin_top = 0
    txBox.margin_bottom = 0
    txBox.margin_left = 0
    txBox.margin_right = 0

    p = tf.paragraphs[0]
    p.text = label_text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(8)
    p.font.name = 'Arial'
    p.font.bold = True
    return pic

# 添加右侧三幅子图
pic1 = add_subimg(sub1_path, 0.3, "(a) Livox Mid-360 LiDAR")
pic2 = add_subimg(sub2_path, 1.7, "(b) Jetson AGX Orin")
pic3 = add_subimg(sub3_path, 3.1, "(c) Ackermann chassis")

# 5. 添加连线（使用 Block Arrow 形状）
def add_arrow(left, top, width):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top),
        Inches(width), Inches(0.08)
    )
    shape.fill.solid()
    # 使用比较好看的科技蓝，如 RGB(0, 114, 198)
    shape.fill.fore_color.rgb = RGBColor(0, 114, 198)
    shape.line.color.rgb = RGBColor(0, 114, 198)

# 使用大概的位置坐标，方便用户到 PPT 中手动拽动微调
add_arrow(2.2, 0.9, 1.3)   # 顶端 -> Livox
add_arrow(1.6, 2.2, 1.9)   # 中间 -> Orin
add_arrow(2.2, 3.4, 1.3)   # 底端 -> Chassis

# 6. 保存输出
out_dir = "/home/sun/phdproject/dqn/DQN8/paperdqn8.3/media"
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "fig_platform.pptx")
prs.save(out_path)
print(f"PPTX successfully saved to {out_path}")
